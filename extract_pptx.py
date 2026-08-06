"""
Extracts text, speaker notes, and table content from .pptx lecture files
into a structured JSON format ready for chunking/embedding in a RAG pipeline.

Usage:
    python extract_pptx.py <input_folder> <output_json> [--course COURSE]

Example:
    python extract_pptx.py data/raw/COMP719/ data/comp719_slides.json --course COMP719

If --course is omitted, the course is inferred from the input folder's name
(e.g. data/raw/COMP719/ -> "COMP719"). Filename-based inference was dropped
since lecture filenames like "L1 - Introduction.pptx" don't reliably encode
a course code -- organize slides into one folder per course instead.
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation


def extract_shape_text(shape):
    """Extract text from a single shape (text box, title, placeholder, etc.)."""
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs)
        if line.strip():
            lines.append(line)
    return "\n".join(lines)


def extract_table_text(shape):
    """Extract text from a table shape as pipe-separated rows."""
    rows_text = []
    table = shape.table
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows_text.append(" | ".join(cells))
    return "\n".join(rows_text)


def extract_slide(slide, slide_num, source_file, course):
    """Pull all text content + metadata from a single slide."""
    text_parts = []
    has_images = False

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = extract_shape_text(shape)
            if text:
                text_parts.append(text)
        elif shape.has_table:
            text_parts.append(extract_table_text(shape))
        elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            has_images = True

    slide_text = "\n".join(text_parts).strip()

    # Speaker notes, if present
    notes_text = ""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = notes_frame.text.strip() if notes_frame else ""

    return {
        "text": slide_text,
        "notes": notes_text,
        "source_file": source_file,
        "slide_number": slide_num,
        "course": course,
        "chunk_id": f"{Path(source_file).stem}_s{slide_num}".lower().replace(" ", "_").replace("-", "_"),
        "has_images": has_images,
    }


def extract_pptx_file(filepath, course):
    """Extract all slides from one .pptx file."""
    prs = Presentation(filepath)
    source_file = Path(filepath).name

    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide(slide, i, source_file, course)
        # Skip completely empty slides (e.g. section dividers with no text)
        if slide_data["text"] or slide_data["notes"]:
            slides_data.append(slide_data)

    return slides_data


def extract_folder(folder_path, output_path, course=None):
    folder = Path(folder_path)
    pptx_files = sorted(folder.glob("*.pptx"))

    if not pptx_files:
        print(f"No .pptx files found in {folder_path}")
        return

    # Fall back to the folder name if --course wasn't passed explicitly
    resolved_course = course or folder.name
    print(f"Using course label: {resolved_course}")

    all_slides = []
    for pptx_file in pptx_files:
        print(f"Extracting: {pptx_file.name}")
        try:
            slides = extract_pptx_file(pptx_file, course=resolved_course)
            all_slides.extend(slides)
            print(f"  -> {len(slides)} slides extracted")
        except Exception as e:
            print(f"  !! Failed to extract {pptx_file.name}: {e}")

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(all_slides, f, indent=2, ensure_ascii=False)

    print(f"\nDone. {len(all_slides)} total slides written to {output_path}")

    # Quick sanity flag worth checking manually: images present with NO text
    # at all -- these slides are currently invisible to text retrieval.
    invisible = [s for s in all_slides if s["has_images"] and not s["text"]]
    if invisible:
        print(f"Note: {len(invisible)} slide(s) have images but no extracted text "
              f"(diagrams/screenshots) -- not searchable by text retrieval yet.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text/notes/tables from .pptx lecture files.")
    parser.add_argument("input_folder", help="Folder containing .pptx files")
    parser.add_argument("output_json", help="Path to write extracted JSON")
    parser.add_argument("--course", help="Course label to tag every slide with "
                                          "(defaults to the input folder's name)")
    args = parser.parse_args()

    extract_folder(args.input_folder, args.output_json, course=args.course)